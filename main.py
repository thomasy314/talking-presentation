import os
import argparse
from pptx import Presentation
import azure.cognitiveservices.speech as speechsdk

OUTPUT_FILE_LOC = "output/"

def get_text_lines_from_file(filename):
    file = open(filename, 'r')
    text_lines = list(filter(lambda a: a != '', file.read().split('\n')))
    file.close()

    return text_lines

def create_presentation_from_notes(per_slide_notes, output_file):
    # Create new presentation
    prs = Presentation()

    # For all lines in per_slide_notes, create new slide and add notes
    for notes in per_slide_notes:
        # create new slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Create slide notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes

    prs.save(f'{output_file}.pptx')

"""
TODO: Convert to batch request instead of synchronous
"""
def create_speech_from_text_list(text_list, output_file):
    # This example requires environment variables named 'SPEECH_KEY' and 'SPEECH_REGION'
    speech_config = speechsdk.SpeechConfig(subscription=os.environ.get('SPEECH_KEY'), region=os.environ.get('SPEECH_REGION'))
    speech_config.speech_synthesis_voice_name='en-US-AvaMultilingualNeural'

    for i, text in enumerate(text_list):
        file_name = f'{output_file}-{i}.wav'
        file_config = speechsdk.audio.AudioOutputConfig(filename=file_name)

        # The neural multilingual voice can speak different languages based on the input text.
        speech_synthesizer = speechsdk.SpeechSynthesizer(speech_config=speech_config, audio_config=file_config)

        speech_synthesis_result = speech_synthesizer.speak_text_async(text).get()

        if speech_synthesis_result.reason == speechsdk.ResultReason.SynthesizingAudioCompleted:
            print('Speech synthesized for text [{}]'.format(file_name))
        elif speech_synthesis_result.reason == speechsdk.ResultReason.Canceled:
            cancellation_details = speech_synthesis_result.cancellation_details
            print('Speech synthesis canceled: {}'.format(cancellation_details.reason))
            if cancellation_details.reason == speechsdk.CancellationReason.Error:
                if cancellation_details.error_details:
                    print('Error details: {}'.format(cancellation_details.error_details))
                    print('Did you set the speech resource key and region values?')

if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--input-file', required=True)
    parser.add_argument('--output-files-prefix', required=True)
    args = parser.parse_args()

    output_file = f'{OUTPUT_FILE_LOC}{args.output_files_prefix}'
    if not os.path.exists(OUTPUT_FILE_LOC):
        os.makedirs(OUTPUT_FILE_LOC)

    text_lines = get_text_lines_from_file(args.input_file)
    create_presentation_from_notes(text_lines, output_file)
    create_speech_from_text_list(text_lines, output_file)
